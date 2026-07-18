from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from copy import deepcopy
import os

# Open template
tmpl_path = '/Users/op/文件/MBA/02 常用文档/9-首都经济贸易大学PPT模板.pptx'
prs = Presentation(tmpl_path)

# Remove all existing slides except layouts
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

L = prs.slide_layouts
COVER = L[0]       # 标题幻灯片
CONTENT = L[1]     # 标题和内容
SECTION = L[2]     # 节标题
TWO_COL = L[3]     # 两栏内容
TITLE_ONLY = L[5]  # 仅标题
BLANK = L[6]       # 空白

def add_text(slide, left, top, width, height, text, size=12, bold=False, color=None, name='Arial', align=None):
    """Add a text box to a slide (EMU values)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = name
    if color:
        p.font.color.rgb = RGBColor(*color)
    if align is not None:
        from pptx.enum.text import PP_ALIGN
        p.alignment = align
    return tf

def add_multiline(slide, x, y, w, h, lines):
    """lines: list of (text, size, bold, color)"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        if color:
            p.font.color.rgb = RGBColor(*color)
        p.space_after = Pt(4)
    return tf

# ===== S1: COVER =====
s = prs.slides.add_slide(COVER)
# The template cover already has the CUEB logo and decorative elements
# Find and update title text
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            if '首都经济贸易大学' in p.text:
                p.text = '首都经济贸易大学  MBA学位论文开题答辩'
            elif '答辩' in p.text and '模板' in p.text:
                p.text = 'REITs驱动的商业地产企业轻资产转型路径与效应研究'
            elif '研究生答辩' in p.text or '博士答辩' in p.text:
                p.text = '王江峰  |  22025023138'
            elif '首经贸' in p.text:
                p.text = '王江峰'
                # Make text bigger
                for run in p.runs:
                    run.font.size = Pt(18)

# ===== S2: FOUR NUMBERS =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '四组数字，一个问题', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
nums = [
    ('77.27%', '大悦城控股 资产负债率'),
    ('70.81%', '新城控股 资产负债率'),
    ('52.2%', '龙湖集团 净负债率'),
    ('14.2%', '华润置地 开发毛利率'),
]
for i, (v, l) in enumerate(nums):
    x = Inches(0.5 + (i%2)*4.8)
    y = Inches(1.2 + (i//2)*1.9)
    # Card background
    sh = s.shapes.add_shape(1, x, y, Inches(4.4), Inches(1.65))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xF5,0xF5,0xF5)
    sh.line.fill.background()
    # Number
    add_text(s, Inches(0.5+(i%2)*4.8+0.3), Inches(1.25+(i//2)*1.9), Inches(3.8), Inches(0.7), v, size=42, bold=True, color=(0x00,0x56,0xA0), name='Arial Black')
    # Label
    add_text(s, Inches(0.5+(i%2)*4.8+0.3), Inches(1.95+(i//2)*1.9), Inches(3.8), Inches(0.35), l, size=11, color=(0x66,0x66,0x66))
add_text(s, Inches(0.6), Inches(5.1), Inches(9), Inches(0.2), '各公司2025年年报  ·  新城同比降2.28pp  ·  华润六年跌36个百分点', size=8, color=(0x99,0x99,0x99))

# ===== S3: CHART SLIDE =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '50.3% → 14.2%', size=28, bold=True, color=(0xCC,0x33,0x33), name='Microsoft YaHei')
# Note: python-pptx can't do charts reliably. Use text + visual bars
years = [('2019','50.3%',50.3),('2020','30.2%',30.2),('2021','24.5%',24.5),('2022','22.2%',22.2),('2023','18.6%',18.6),('2024','14.2%',14.2),('2025','15.5%',15.5)]
for i, (yr, val, vnum) in enumerate(years):
    x = Inches(0.5 + i*1.35)
    h = Inches(vnum/55 * 3.5)
    bottom = Inches(4.3)
    sh = s.shapes.add_shape(1, x, bottom - h, Inches(1.15), h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x00,0x56,0xA0) if i < 6 else RGBColor(0xCC,0x33,0x33)
    sh.line.fill.background()
    add_text(s, x, Inches(4.35), Inches(1.15), Inches(0.2), val, size=9, bold=True, color=(0x00,0x56,0xA0))
    add_text(s, x, Inches(4.55), Inches(1.15), Inches(0.2), yr, size=8, color=(0x99,0x99,0x99))
# Right side text
add_text(s, Inches(7.2), Inches(1.2), Inches(2.5), Inches(3.0), '2020年末\n投资性房地产\n2,363亿元\n\n占总资产超30%\n\n"资产沉淀→\n负债攀升→\n盈利承压"', size=10, color=(0x33,0x33,0x33))
add_text(s, Inches(0.6), Inches(5.1), Inches(9), Inches(0.2), '华润置地年度报告（2019-2025）', size=8, color=(0x99,0x99,0x99))

# ===== S4: TURNING POINT =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '转机正在发生', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
cards = [('180.22亿','华润万象生活','2025年营收'),('1,250亿','凯德集团','FUM（新元）'),('51.8%','经常性利润占比','首次过半')]
for i, (n, l1, l2) in enumerate(cards):
    x = Inches(0.5 + i*3.2)
    sh = s.shapes.add_shape(1, x, Inches(1.2), Inches(2.9), Inches(2.1))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xF0,0xF5,0xFA)
    sh.line.fill.background()
    add_text(s, Inches(0.5+i*3.2+0.2), Inches(1.3), Inches(2.5), Inches(0.7), n, size=30, bold=True, color=(0x00,0x56,0xA0), name='Arial Black')
    add_text(s, Inches(0.5+i*3.2+0.2), Inches(2.1), Inches(2.5), Inches(0.3), l1, size=10, bold=True, color=(0x33,0x33,0x33))
    add_text(s, Inches(0.5+i*3.2+0.2), Inches(2.4), Inches(2.5), Inches(0.25), l2, size=9, color=(0x99,0x99,0x99))
# Bottom bar
sh = s.shapes.add_shape(1, Inches(0.5), Inches(3.6), Inches(9.0), Inches(1.2))
sh.fill.solid()
sh.fill.fore_color.rgb = RGBColor(0x00,0x56,0xA0)
sh.line.fill.background()
add_text(s, Inches(0.8), Inches(3.75), Inches(8.4), Inches(0.5), 'REITs 如何驱动企业完成这场转型？', size=16, bold=True, color=(0xFF,0xFF,0xFF), name='Microsoft YaHei')
add_text(s, Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.35), '路径 · 效应 · 条件 —— 三个维度，三家企业，一个答案', size=11, color=(0xCC,0xDD,0xFF))

# ===== S5: THREE CASES =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '三家企业，三条路径', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
cases = [
    ('华润置地+万象生活','主案例','双平台：商管分拆+REITs发行\n港交所 2020-2025\n国内首创模式',RGBColor(0x00,0x56,0xA0)),
    ('凯德集团','验证案例','私募基金+REITs 资本闭环\n新交所 2020-2025\n20年成熟经验',RGBColor(0x00,0x96,0x88)),
    ('苏宁易购','反面案例','类REITs/ABS  65%→91%\n深交所 2015-2023\n出表不降杠杆',RGBColor(0xCC,0x33,0x33)),
]
for i, (name, role, desc, clr) in enumerate(cases):
    x = Inches(0.5 + i*3.2)
    sh = s.shapes.add_shape(1, x, Inches(1.2), Inches(2.9), Inches(2.6))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xF5,0xF5,0xF5)
    sh.line.fill.background()
    # Color accent bar on top
    bar_sh = s.shapes.add_shape(1, x, Inches(1.2), Inches(2.9), Inches(0.05))
    bar_sh.fill.solid()
    bar_sh.fill.fore_color.rgb = clr
    bar_sh.line.fill.background()
    add_text(s, Inches(0.5+i*3.2+0.2), Inches(1.35), Inches(2.5), Inches(0.4), name, size=14, bold=True, color=(0x33,0x33,0x33), name='Microsoft YaHei')
    add_text(s, Inches(0.5+i*3.2+0.2), Inches(1.75), Inches(2.5), Inches(0.25), role, size=9, bold=True, color=(clr[0],clr[1],clr[2]))
    add_text(s, Inches(0.5+i*3.2+0.2), Inches(2.15), Inches(2.5), Inches(1.4), desc, size=10, color=(0x66,0x66,0x66))
add_text(s, Inches(0.6), Inches(4.4), Inches(9), Inches(0.3), '正向验证 + 反向证伪 + 跨国对比 = 完整证据链', size=11, bold=True, color=(0x00,0x56,0xA0))

# ===== S6: LITERATURE MAP =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '本文在文献坐标系中的位置', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
# Quadrants
ox, oy, cw, ch = Inches(0.4), Inches(1.2), Inches(2.3), Inches(1.9)
fills = [RGBColor(0xF5,0xF5,0xF5), RGBColor(0xEE,0xF2,0xF8), RGBColor(0xF5,0xF5,0xF5), RGBColor(0xEE,0xF2,0xF8)]
for r in range(2):
    for c in range(2):
        sh = s.shapes.add_shape(1, ox + c*cw, oy + r*ch, cw, ch)
        sh.fill.solid(); sh.fill.fore_color.rgb = fills[r*2+c]; sh.line.fill.background()
# Cross lines
s.shapes.add_shape(1, ox + cw, oy, Pt(1), ch*2).fill.solid()
s.shapes.add_shape(1, ox, oy + ch, cw*2, Pt(1)).fill.solid()
# Axis labels
add_text(s, Inches(0.4), Inches(1.02), Inches(2.3), Inches(0.15), '金融工具视角', size=7, color=(0x99,0x99,0x99))
add_text(s, Inches(2.7), Inches(1.02), Inches(2.3), Inches(0.15), '企业战略转型视角', size=7, bold=True, color=(0x00,0x56,0xA0))
# Dots (gray) + our paper (red)
dots = [(0.12,0.55),(0.22,0.7),(0.08,0.82),(0.55,0.6),(0.58,0.75),(0.48,0.85),(0.18,0.15),(0.25,0.28),(0.6,0.22),(0.72,0.5)]
for dx, dy in dots:
    sh = s.shapes.add_shape(9, int(ox + dx*cw*2), int(oy + dy*ch*2), Pt(10), Pt(10))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xBB,0xBB,0xBB); sh.line.fill.background()
# Our paper
sh = s.shapes.add_shape(9, int(ox + 0.78*cw*2), int(oy + 0.12*ch*2), Pt(18), Pt(18))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xCC,0x33,0x33); sh.line.fill.background()
add_text(s, Inches(0.4+0.78*2.3*2-0.05), Inches(1.2+0.12*1.9*2-0.15), Inches(0.3), Inches(0.2), '本文', size=6, bold=True, color=(0xFF,0xFF,0xFF))
# Right text
add_text(s, Inches(5.5), Inches(1.3), Inches(4.0), Inches(3.0),
    '定位：企业战略转型视角\n× 多维度系统评估\n\n构建"交易结构-出表效果\n  -转型效应"统一分析框架\n\n五维度财务穿透\n+ 三案例双向验证', size=10, color=(0x33,0x33,0x33))

# ===== S7: THREE GAPS =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '三个研究不足', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
gaps = [
    ('01','交易结构→出表效果→转型效应：尚未整合','Hardin/Wu · 原野/彭 · 王宇/马——分属金融学、会计学、公司财务，未形成统一框架'),
    ('02','转型评估停留在单一指标','Sohn→ROE · Hardin→负债率 · 吴晓波→盈利波动——各从一个维度，缺少系统工具'),
    ('03','正反案例双向验证缺失','多单案例描述(万达/越秀/凯德)，华润"双平台"首创实践文献空白'),
]
for i, (n, t, d) in enumerate(gaps):
    y = Inches(1.3 + i*1.35)
    sh = s.shapes.add_shape(1, Inches(0.5), y, Inches(9.0), Inches(1.1))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF5,0xF5,0xF5); sh.line.fill.background()
    # Left accent bar
    ab = s.shapes.add_shape(1, Inches(0.5), y, Inches(0.06), Inches(1.1))
    ab.fill.solid(); ab.fill.fore_color.rgb = RGBColor(0x00,0x56,0xA0); ab.line.fill.background()
    add_text(s, Inches(0.8), Inches(1.3+i*1.35+0.05), Inches(0.4), Inches(0.3), n, size=16, bold=True, color=(0x00,0x56,0xA0), name='Arial Black')
    add_text(s, Inches(1.3), Inches(1.3+i*1.35+0.05), Inches(7.8), Inches(0.3), t, size=12, bold=True, color=(0x33,0x33,0x33))
    add_text(s, Inches(1.3), Inches(1.3+i*1.35+0.45), Inches(7.8), Inches(0.5), d, size=9, color=(0x66,0x66,0x66))

# ===== S8: HYPOTHESES =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '五个研究假设', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
hyps = [
    'H₁：REITs发行 → 经常性收入占比显著提高',
    'H₂：轻资产转型 → 净负债率显著下降',
    'H₃：运营能力是转型成功的必要条件',
    'H₄：真实出售的效应 远大于 售后租回',
    'H₅：完整资本循环 → 转型效应更持久',
]
for i, h in enumerate(hyps):
    y = Inches(1.2 + i*0.68)
    sh = s.shapes.add_shape(1, Inches(0.5), y, Inches(9.0), Inches(0.52))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF0,0xF5,0xFA); sh.line.fill.background()
    add_text(s, Inches(0.75), Inches(1.2+i*0.68+0.06), Inches(0.35), Inches(0.3), '0'+str(i+1), size=13, bold=True, color=(0x00,0x56,0xA0), name='Arial Black')
    add_text(s, Inches(1.2), Inches(1.2+i*0.68+0.06), Inches(8.0), Inches(0.3), h, size=13, color=(0x33,0x33,0x33))
add_text(s, Inches(0.6), Inches(4.8), Inches(9), Inches(0.25), '理论依据：Hardin & Wu (2010) · Sohn et al. (2013) · 原野、彭晓松 (2019) · 刘月龙 (2023)', size=8, color=(0x99,0x99,0x99))

# ===== S9: FRAMEWORK =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(9), Inches(0.7), '分析框架', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
layers = [
    ('REITs驱动机制','制度功能 × 企业战略需求'),
    ('路径选择','华润双平台 vs 凯德资本闭环'),
    ('效应评估','财务效应(5维) + 运营效应(3维)'),
]
for i, (t, d) in enumerate(layers):
    y = Inches(1.2 + i*1.1)
    sh = s.shapes.add_shape(1, Inches(0.5), y, Inches(4.5), Inches(0.9))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF0,0xF5,0xFA); sh.line.fill.background()
    ab = s.shapes.add_shape(1, Inches(0.5), y, Inches(0.06), Inches(0.9))
    ab.fill.solid(); ab.fill.fore_color.rgb = RGBColor(0x00,0x56,0xA0); ab.line.fill.background()
    add_text(s, Inches(0.75), Inches(1.2+i*1.1+0.05), Inches(1.5), Inches(0.3), t, size=12, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
    add_text(s, Inches(2.3), Inches(1.2+i*1.1+0.1), Inches(2.5), Inches(0.6), d, size=10, color=(0x33,0x33,0x33))
    if i < 2:
        add_text(s, Inches(2.5), Inches(1.2+i*1.1+0.7), Inches(0.5), Inches(0.2), '↓', size=13, bold=True, color=(0x00,0x56,0xA0))
# Right side
add_text(s, Inches(5.5), Inches(1.2), Inches(4.2), Inches(0.3), '五维度财务穿透', size=12, bold=True, color=(0x00,0x56,0xA0))
dims = ['① 资产轻量化','② 收入结构','③ 盈利质量','④ 现金流','⑤ 资本效率']
for i, d in enumerate(dims):
    add_text(s, Inches(5.5), Inches(1.6+i*0.28), Inches(4.2), Inches(0.22), d, size=10, color=(0x33,0x33,0x33))
add_text(s, Inches(5.5), Inches(3.3), Inches(4.2), Inches(0.3), '五阶段转型路径', size=12, bold=True, color=(0x00,0x56,0xA0))
add_text(s, Inches(5.5), Inches(3.7), Inches(4.2), Inches(0.3), '持有 → 剥离 → 平台 → REITs → 循环', size=10, bold=True, color=(0x00,0x56,0xA0))
add_text(s, Inches(0.6), Inches(5.1), Inches(9), Inches(0.2), '数据来源：年报 + REITs文件 + Wind/戴德梁行/彭博', size=8, color=(0x99,0x99,0x99))

# ===== S10: INNOVATIONS =====
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.6), Inches(0.3), Inches(5), Inches(0.7), '四个创新点', size=28, bold=True, color=(0x00,0x56,0xA0), name='Microsoft YaHei')
innov = [
    ('①','首次系统研究"双平台模式"','华润"商管分拆+REITs发行"国内首创，文献空白'),
    ('②','构建五维度财务穿透框架','超越ROE/负债率，资产→收入→盈利→现金流→资本效率'),
    ('③','引入反面案例双向验证','苏宁"出表不降杠杆"↔华润/凯德成功，正反对照'),
    ('④','跨国对比增强普适性','华润(央企)↔凯德(新加坡)，制度环境差异检验'),
]
for i, (n, t, d) in enumerate(innov):
    y = Inches(1.1 + i*0.8)
    sh = s.shapes.add_shape(1, Inches(0.5), y, Inches(6.2), Inches(0.65))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF0,0xF5,0xFA); sh.line.fill.background()
    ab = s.shapes.add_shape(1, Inches(0.5), y, Inches(0.05), Inches(0.65))
    ab.fill.solid(); ab.fill.fore_color.rgb = RGBColor(0x00,0x56,0xA0); ab.line.fill.background()
    add_text(s, Inches(0.75), Inches(1.1+i*0.8+0.05), Inches(0.3), Inches(0.25), n, size=14, bold=True, color=(0x00,0x56,0xA0), name='Arial Black')
    add_text(s, Inches(1.15), Inches(1.1+i*0.8+0.03), Inches(5.3), Inches(0.25), t, size=11, bold=True, color=(0x33,0x33,0x33))
    add_text(s, Inches(1.15), Inches(1.1+i*0.8+0.3), Inches(5.3), Inches(0.2), d, size=8, color=(0x66,0x66,0x66))
# Outline
add_text(s, Inches(7.0), Inches(1.1), Inches(2.8), Inches(0.3), '七章结构', size=11, bold=True, color=(0x00,0x56,0xA0))
chapters = ['绪论','文献综述','研究设计','华润分析','凯德分析','苏宁反面','结论建议']
for i, ch in enumerate(chapters):
    add_text(s, Inches(7.0), Inches(1.5+i*0.28), Inches(2.8), Inches(0.22), f'第{i+1}章  {ch}', size=9, color=(0x66,0x66,0x66))
# Timeline
add_text(s, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.25), '2026.05-07 文献 → 08-09 资料 → 10-11 华润 → 11-12 凯德 → 2027.01-02 初稿 → 02-03 修改 → 03-04 定稿', size=9, color=(0x99,0x99,0x99))

# ===== S11: THANK YOU =====
s = prs.slides.add_slide(COVER)
for sh in s.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            if '首都经济贸易大学' in p.text:
                p.text = ''
            elif '答辩' in p.text and '模板' in p.text:
                p.text = '感谢各位老师'
            elif '研究生答辩' in p.text or '博士答辩' in p.text:
                p.text = '恳请批评指正'
            elif '首经贸' in p.text:
                p.text = ''

# Save
out_path = '/Users/op/WorkBuddy/论文优化/开题答辩PPT.pptx'
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Slides: {len(prs.slides)}')
