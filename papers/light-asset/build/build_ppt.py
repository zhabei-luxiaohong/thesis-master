# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.6)

C  = RGBColor(0x8B,0x00,0x00)
D  = RGBColor(0x1A,0x1A,0x2E)
G  = RGBColor(0x99,0x99,0x99)
W  = RGBColor(0xFF,0xFF,0xFF)
LG = RGBColor(0xF2,0xF2,0xF2)
P  = RGBColor(0xFF,0xF0,0xF0)

A_L = PP_ALIGN.LEFT
A_C = PP_ALIGN.CENTER
A_R = PP_ALIGN.RIGHT

LQ = '\u201c'  # left Chinese quote
RQ = '\u201d'  # right Chinese quote

def bar(s):
    sh = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.04))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C
    sh.line.fill.background()

def T(s, x, y, w, h, t, sz=12, cl=D, b=False, al=None, fn='Arial'):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = t
    p.font.size = Pt(sz)
    p.font.color.rgb = cl
    p.font.bold = b
    p.font.name = fn
    if al is not None:
        p.alignment = al
    return tf

def tag(s, t):
    T(s, 0.6, 0.88, 5, 0.2, t, 8, C, True)

def tit(s, t):
    T(s, 0.6, 0.18, 8.8, 0.6, t, 22, D, True, fn='Arial Black')

def src(s, t):
    T(s, 0.6, 5.2, 8.8, 0.18, t, 6, G)

def kpi(s, x, y, n, l):
    T(s, x, y, 2.5, 0.5, n, 18, C, True, A_C, 'Arial Black')
    T(s, x, y+0.48, 2.5, 0.4, l, 8, D, al=A_C)

def box(s, x, y, w, h, fill=P):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()

# ==============================================================
# S1: COVER
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = D

t = 'REITs\u9a71\u52a8\u7684\u5546\u4e1a\u5730\u4ea7\u4f01\u4e1a\n\u8f7b\u8d44\u4ea7\u8f6c\u578b\u8def\u5f84\u4e0e\u6548\u5e94\u7814\u7a76'
T(s, 1, 1.2, 8, 1.2, t, 30, W, True, A_C, 'Arial Black')
T(s, 1, 2.6, 8, 0.4, 'MBA学位论文开题答辩', 14, RGBColor(0xCC,0x66,0x77), al=A_C)
T(s, 1, 4.0, 8, 0.3, '王江峰  22025023138    |    首都经济贸易大学 MBA教育中心    |    2026年5月', 10, G, al=A_C)

# ==============================================================
# S2: PEST
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, 'PEST\uff1a\u56db\u91cd\u529b\u91cf\u91cd\u5851\u5546\u4e1a\u5730\u4ea7')
tag(s, 'PART 1  \u4e3a\u4ec0\u4e48\u9009\u8fd9\u4e2a\u9898')

pest = [
    ('P', '政策驱动', 'REITs试点\u2192消费扩容\n注册制落地\uff0c2,100亿/58只', '\u2192 退出通道打开'),
    ('E', '融资收紧', '开发毛利率50.3%\u219214.2%\n重资产模式难以为继', '\u2192 被迫转向轻资产'),
    ('S', '消费分化', '电商冲击+消费升级并存\n优质购物中心出租率>95%', '\u2192 运营能力成核心'),
    ('T', '管理升级', '数字化运营平台成熟\n标准化体系可复制', '\u2192 轻资产条件具备'),
]
for i, (p, t, d, e) in enumerate(pest):
    x, y = 0.5 + (i%2)*4.7, 1.3 + (i//2)*2.0
    box(s, x, y, 4.4, 1.8, LG)
    T(s, x+0.15, y+0.08, 0.4, 0.4, p, 20, C, True, fn='Arial Black')
    T(s, x+0.65, y+0.08, 3.5, 0.25, t, 13, D, True)
    T(s, x+0.25, y+0.5, 3.9, 0.9, d, 9, G)
    T(s, x+0.25, y+1.45, 3.9, 0.25, e, 9, C, True)

# ==============================================================
# S3: 行业困境
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '四重压力叠加：行业面临结构性困局')
tag(s, 'PART 1  为什么选这个题')
kpi(s, 0.7, 1.6, '77.27%', '大悦城控股 资产负债率')
kpi(s, 3.5, 1.6, '70.81%', '新城控股 资产负债率')
kpi(s, 6.3, 1.6, '52.2%', '龙湖集团 净负债率')
box(s, 0.6, 2.8, 8.8, 1.0, P)
T(s, 0.8, 2.9, 8.4, 0.8, '融资渠道收窄 + 销售回款放缓 + 存量资产难以盘活\n头部杠杆率较2024年小幅改善（新城降2.28pp），但行业困境未解。\n华润置地开发毛利率六年跌36个百分点：50.3% \u2192 14.2%（2024），2025触底回升至15.5%。', 10, D)
src(s, '各公司2025年年报; Wind数据')

# ==============================================================
# S4: 华润数据条
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '华润置地：六年数据轨迹')
tag(s, 'PART 1  为什么选这个题')

years = [
    ('2019', '50.3%'), ('2020', '30.2%'), ('2021', '24.5%'),
    ('2022', '22.2%'), ('2023', '18.6%'), ('2024', '14.2%'), ('2025', '15.5%')
]
T(s, 0.6, 1.3, 3, 0.25, '开发业务毛利率（%）', 11, D, True)
for i, (yr, val) in enumerate(years):
    x = 0.6 + i*1.3
    v = float(val.replace('%', ''))
    bar_h = v/50 * 8 * 0.055
    box(s, x, 1.7, 1.15, 0.08, C)
    box(s, x, 1.78, 1.15, bar_h, C)
    T(s, x, 1.8+bar_h, 1.15, 0.2, val, 10, C, True, A_C)
    T(s, x, 1.55, 1.15, 0.2, yr, 9, G, al=A_C)

text_s4 = (
    '2020年末投资性房地产 2,363亿元，占总资产超过30%。\n'
    + LQ + '资产沉淀 \u2192 负债攀升 \u2192 盈利承压' + RQ
    + '恶性循环，迫使企业寻求转型。\n'
    + '2024年经常性利润占比首次过半（51.8%），轻资产转型拐点已至。'
)
T(s, 0.6, 3.7, 8.8, 1.0, text_s4, 10, D)
src(s, '华润置地年度报告（2019-2025）')

# ==============================================================
# S5: 三个研究问题
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '转型正在发生\u2014\u2014但三个问题待解')
tag(s, 'PART 1  为什么选这个题')
kpi(s, 0.7, 1.4, '78%\u219242%', '凯德集团负债率（2001\u21922021）')
kpi(s, 3.5, 1.4, '180.22亿', '华润万象生活营收（2025）')
kpi(s, 6.3, 1.4, '51.8%', '经常性利润占比首过半（2024）')

box(s, 0.6, 2.6, 8.8, 2.4, P)
T(s, 0.8, 2.7, 8.4, 0.3, '核心研究问题', 13, C, True)
q_text = (
    'Q1  路径：REITs如何驱动轻资产转型？关键环节和组织变革是什么？\n'
    + 'Q2  效应：如何系统量化财务与运营效应？对三张报表的具体影响？\n'
    + 'Q3  条件：成功的边界条件是什么？如何避免' + LQ + '出表不降杠杆' + RQ + '？\n\n'
    + '\u2192 以华润（主案例）+凯德（验证）+苏宁（反面）三案例对比设计回答'
)
T(s, 0.8, 3.1, 8.4, 1.7, q_text, 10, D)
T(s, 0.7, 5.15, 8.6, 0.2, '国际对标：西蒙\u00b7凯德\u00b7西田  |  国内：2025注册制\u00b72,100亿/58只\u00b7消费REITs扩容', 7, G)

# ==============================================================
# S6: 文献坐标系
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '本文在文献坐标系中的位置')
tag(s, 'PART 2  文献定位')

ox, oy, gw, gh = 0.5, 1.3, 5.0, 3.8
for r in range(2):
    for c in range(2):
        clrs = [[LG, P], [RGBColor(0xF0,0xF0,0xF0), RGBColor(0xF5,0xF5,0xF5)]]
        box(s, ox+c*gw/2, oy+r*gh/2, gw/2, gh/2, clrs[r][c])

# Cross lines
s.shapes.add_shape(1, Inches(ox+gw/2), Inches(oy+0.1), Inches(0.008), Inches(gh-0.2)).fill.solid()
s.shapes.add_shape(1, Inches(ox+0.1), Inches(oy+gh/2), Inches(gw-0.2), Inches(0.008)).fill.solid()

T(s, ox, oy-0.05, gw/2, 0.18, '金融工具视角', 8, G, al=A_C)
T(s, ox+gw/2, oy-0.05, gw/2, 0.18, '企业战略转型视角', 8, C, True, A_C)
T(s, ox-0.38, oy+0.5, 0.28, 2.5, '多维度评估\n\n\u2191\n\n\n\u2193\n\n单一指标', 7, C, True, A_C)

dots = [(0.15,0.65),(0.2,0.82),(0.1,0.92),(0.55,0.7),(0.6,0.83),(0.5,0.93),(0.2,0.15),(0.3,0.3)]
for dx, dy in dots:
    dt = s.shapes.add_shape(9, Inches(ox+dx*gw), Inches(oy+dy*gh), Inches(0.2), Inches(0.2))
    dt.fill.solid()
    dt.fill.fore_color.rgb = G
    dt.line.fill.background()

dt = s.shapes.add_shape(9, Inches(ox+0.75*gw), Inches(oy+0.15*gh), Inches(0.35), Inches(0.35))
dt.fill.solid()
dt.fill.fore_color.rgb = C
dt.line.color.rgb = C
dt.line.width = Pt(2)
T(s, ox+0.75*gw, oy+0.13*gh, 0.35, 0.3, '本文', 7, W, True, A_C, 'Arial Black')

T(s, 5.8, 1.4, 3.8, 0.3, '本文定位：右上角', 13, C, True)
T(s, 5.8, 1.8, 3.8, 3.0, '视角：企业战略转型\n  超越金融工具\n\n深度：五维度系统评估\n  财务穿透框架\n\n设计：三案例双向验证\n  华润+凯德+苏宁正反对照\n\n\u2192 右上角 = 学术贡献', 9, D)

# ==============================================================
# S7: 三个缺口
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '三个缺口与本文切入点')
tag(s, 'PART 2  文献定位')

gaps = [
    ('01  交易结构\u2192出表效果\u2192转型效应：尚未整合',
     'Hardin/Wu + 原野/彭 + 王宇/马 \u2192 分属金融学、会计学、公司财务，未形成从交易到转型的统一分析框架'),
    ('02  转型评估停留在单一指标',
     'Sohn\u2192ROE / Hardin\u2192负债率 / 吴晓波\u2192盈利波动 \u2192 各从一个维度切入，缺少系统的多维度评估工具'),
    ('03  正反案例双向验证缺失',
     '现有多单案例描述（万达/越秀/凯德），华润' + LQ + '双平台' + RQ + '模式国内首创，文献几近空白'),
]
for i, (num, desc) in enumerate(gaps):
    y = 1.4 + i*1.3
    box(s, 0.6, y, 8.8, 1.15, P)
    T(s, 0.8, y+0.05, 8.4, 0.25, num, 12, D, True)
    T(s, 0.8, y+0.35, 8.4, 0.7, desc, 8, G)

T(s, 0.6, 5.2, 8.8, 0.2,
  '\u2192 本文定位：企业战略转型视角，构建' + LQ + '交易结构-出表效果-转型效应' + RQ + '分析框架，三案例正反验证，五维度系统评估',
  7, C, True)

# ==============================================================
# S8: 研究设计
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '研究设计总览')
tag(s, 'PART 3  如何研究')

T(s, 0.6, 1.3, 4.5, 0.25, '核心问题 \u2192 五个假设', 12, C, True)
box(s, 0.6, 1.7, 4.5, 2.4, LG)
h_text = (
    '核心：REITs如何驱动轻资产转型？\n\n'
    'H1\uff1aREITs  \u2192 经常性收入占比\u2191\n'
    'H2\uff1a转型  \u2192 净负债率\u2193\n'
    'H3\uff1a运营能力 = 转型必要条件\n'
    'H4\uff1a真实出售 >> 售后租回\n'
    'H5\uff1a完整资本循环 \u2192 效应更持久\n\n'
    '理论依据：Hardin/Wu(2010) \u00b7 Sohn(2013)\n'
    '                    原野/彭晓松(2019) \u00b7 刘月龙(2023)'
)
T(s, 0.8, 1.8, 4.1, 2.2, h_text, 9, D)

T(s, 5.3, 1.3, 4.3, 0.25, '三案例对比设计', 12, C, True)
cases = [
    ('华润置地/万象生活', '主案例\u00b7深度分析', '双平台：分拆+REITs', '港交所 2020-2025', C),
    ('凯德集团', '验证\u00b7跨国对比', '私募基金+REITs闭环', '新交所 2020-2025', RGBColor(0x4A,0x4A,0x4A)),
    ('苏宁易购', '反面\u00b7失败反证', '类REITs/ABS\n65%\u219291%', '深交所 2015-2023', RGBColor(0xB5,0x35,0x35)),
]
for i, (name, role, mode, data, clr) in enumerate(cases):
    y = 1.75 + i*1.0
    box(s, 5.3, y, 4.3, 0.85, LG)
    T(s, 5.45, y+0.02, 2.0, 0.2, name, 12, D, True)
    T(s, 7.5, y+0.02, 1.9, 0.2, role, 8, clr, True)
    T(s, 5.45, y+0.3, 2.5, 0.22, mode, 8, D)
    T(s, 7.5, y+0.3, 1.9, 0.22, data, 7, G)
T(s, 5.3, 5.0, 4.3, 0.18, 'Yin(2018)案例选择 | 理论抽样 | 七选三', 7, G)

# ==============================================================
# S9: 分析框架
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '分析框架')
tag(s, 'PART 3  如何研究')

T(s, 0.6, 1.3, 4.5, 0.25, '三层分析架构', 12, C, True)
for i, (t, d) in enumerate([
    ('REITs驱动机制', '制度功能 \u00d7 企业战略需求'),
    ('路径选择', '华润双平台 vs 凯德资本闭环'),
    ('效应评估', '财务效应(5维) + 运营效应(3维)'),
]):
    y = 1.7 + i*1.0
    box(s, 0.6, y, 4.5, 0.82, LG)
    T(s, 0.8, y+0.05, 1.4, 0.25, t, 12, C, True, fn='Arial Black')
    T(s, 2.3, y+0.05, 2.6, 0.6, d, 10, D)
    if i < 2:
        T(s, 2.6, y+0.7, 0.5, 0.2, '\u2193', 14, C, True, A_C)

T(s, 5.5, 1.3, 4.0, 0.25, '五维度财务穿透框架', 12, C, True)
box(s, 5.5, 1.7, 4.0, 1.6, LG)
dims = [
    '\u2460 资产轻量化  投资性房地产/总资产',
    '\u2461 收入结构    经常性收入/总收入',
    '\u2462 盈利质量    NOI利润率\u00b7经常性利润占比',
    '\u2463 现金流      经营现金流/带息债务',
    '\u2464 资本效率    资本化率\u00b7ROIC',
]
for i, d in enumerate(dims):
    T(s, 5.7, 1.8+i*0.28, 3.5, 0.22, d, 9, D)

T(s, 5.5, 3.5, 4.0, 0.25, '五阶段转型路径模型', 12, C, True)
T(s, 5.5, 3.85, 4.0, 1.0, '重资产持有  \u2192  剥离分拆  \u2192  平台建设  \u2192  REITs退出  \u2192  资本循环', 9, C, True)
T(s, 0.6, 5.2, 8.8, 0.18, '方法：文献研究+多案例(Yin,2018)+财务穿透+跨国对比 | 数据：年报+REITs文件+Wind/戴德梁行/彭博', 6, G)

# ==============================================================
# S10: 创新点 + 大纲
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = W
bar(s)
tit(s, '创新点与论文大纲')
tag(s, 'PART 4  总结')

T(s, 0.6, 1.3, 6.0, 0.25, '四个创新点', 12, C, True)
innov = [
    ('\u2460 首次系统研究' + LQ + '双平台模式' + RQ, '华润' + LQ + '商管分拆+REITs发行' + RQ + '国内首创，文献几近空白'),
    ('\u2461 构建五维度财务穿透分析框架', '超越ROE/负债率，资产\u2192收入\u2192盈利\u2192现金流\u2192资本效率'),
    ('\u2462 引入反面案例双向验证机制', '苏宁' + LQ + '出表不降杠杆' + RQ + '\u2194华润/凯德成功，正反对照'),
    ('\u2463 跨国对比增强结论普适性', '华润(央企)\u2194凯德(新加坡)，制度环境差异检验'),
]
for i, (n, d) in enumerate(innov):
    y = 1.7 + i*0.72
    T(s, 0.6, y, 6.0, 0.25, n, 11, D, True)
    T(s, 0.8, y+0.25, 5.5, 0.3, d, 8, G)

T(s, 7.0, 1.3, 2.5, 0.25, '七章结构', 12, C, True)
for i, ch in enumerate(['绪论', '文献综述', '研究设计', '华润分析', '凯德分析', '苏宁反面', '结论建议']):
    T(s, 7.0, 1.7+i*0.4, 2.5, 0.35, '\u7b2c' + str(i+1) + '\u7ae0  ' + ch, 10, D, True)

T(s, 0.6, 5.2, 8.8, 0.18,
  '进度：2026.05-07文献 \u2192 08-09资料 \u2192 10-11华润 \u2192 11-12凯德 \u2192 2027.01-02初稿 \u2192 02-03修改 \u2192 03-04定稿',
  7, G)

# ==============================================================
# S11: 致谢
# ==============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = D
bar(s)
T(s, 1, 2.2, 8, 1, '感谢各位老师', 32, W, True, A_C, 'Arial Black')
T(s, 1, 3.5, 8, 0.5, '恳请批评指正', 14, G, al=A_C)

prs.save('/Users/op/WorkBuddy/论文优化/开题答辩PPT.pptx')
print('Done: 11 slides')
